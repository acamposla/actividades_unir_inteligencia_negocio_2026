"""
Crear plantilla PowerPoint con tema profesional para Quarto.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import os

def set_background_color(slide, color_hex):
    """Establecer color de fondo del slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex)

def create_template():
    """Crear plantilla de PowerPoint con tema profesional."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colores del tema (azul corporativo)
    PRIMARY_COLOR = "1E3A5F"      # Azul oscuro
    SECONDARY_COLOR = "2E86AB"    # Azul medio
    ACCENT_COLOR = "F18F01"       # Naranja
    LIGHT_BG = "F5F7FA"           # Gris muy claro

    # Layout 0: Title Slide (Portada)
    title_layout = prs.slide_layouts[0]

    # Layout 1: Title and Content
    content_layout = prs.slide_layouts[1]

    # Layout 5: Title Only (para secciones)
    section_layout = prs.slide_layouts[5]

    # Crear slides de muestra para definir estilos

    # Slide 1: Portada de ejemplo
    slide1 = prs.slides.add_slide(title_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "Título de la Presentación"
    subtitle.text = "Subtítulo | Autor | Fecha"

    # Aplicar formato al título
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor.from_string(PRIMARY_COLOR)

    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor.from_string(SECONDARY_COLOR)

    # Slide 2: Contenido de ejemplo
    slide2 = prs.slides.add_slide(content_layout)
    title2 = slide2.shapes.title
    body2 = slide2.placeholders[1]

    title2.text = "Título del Slide"
    body2.text = "Contenido de ejemplo"

    for paragraph in title2.text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor.from_string(PRIMARY_COLOR)

    for paragraph in body2.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor.from_string("333333")

    # Guardar plantilla
    output_path = os.path.join(os.path.dirname(__file__), 'plantilla_tema.pptx')
    prs.save(output_path)
    print(f"Plantilla creada: {output_path}")
    return output_path

if __name__ == '__main__':
    create_template()
